import streamlit as st
import google.generativeai as genai
import PyPDF2
import docx
import os
import mysql.connector
from mysql.connector import Error
from dotenv import load_dotenv
from werkzeug.security import generate_password_hash, check_password_hash
import pandas as pd
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import nltk
from nltk.sentiment import SentimentIntensityAnalyzer
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import re
import streamlit as st  
import pandas as pd  
import matplotlib.pyplot as plt
from mpl_toolkits.mplot3d import Axes3D  
import seaborn as sns  
import googletrans
import numpy as np 
import io
import datetime
from datetime import datetime, timedelta
from fpdf import FPDF  
from pptx import Presentation 
def trainer_dashboard():  
   st.header("Trainer Dashboard")  
   tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs(["Upload Curriculum", "Generate Question Bank", "View Questions", "Review Feedback", "Clear History", "Download Questions"])  
  
   with tab1:  
      st.subheader("Upload Curriculum")  
      technology = st.text_input("Technology", key="upload_technology")  
      topics = st.text_area("Topics (one per line)", key="upload_topics")  
      uploaded_file = st.file_uploader("Upload curriculum file", type=['txt', 'pdf', 'docx'], key="curriculum_file")  
  
      if st.button("Upload Curriculum", key="upload_curriculum_button"):  
        topic_list = [topic.strip() for topic in topics.split('\n') if topic.strip()]  
        content = ""  
  
        if uploaded_file is not None:  
           try:  
              file_content = extract_text_from_file(uploaded_file)  
              content = file_content  
              topic_list.extend([topic.strip() for topic in file_content.split('\n') if topic.strip()])  
              topic_list = list(set(topic_list)) # Remove duplicates  
           except ValueError as e:  
              st.error(f"Error processing file: {str(e)}")  
              return  
  
        if upload_curriculum(technology, topic_list, content):  
           st.success("Curriculum uploaded successfully!")  
        else:  
           st.error("Failed to upload curriculum")  
  
   
   with tab2:
    st.subheader("Generate Question Bank")
    
    curricula = get_all_curricula()
    
    if not curricula:
        st.warning("No curricula available. Please upload a curriculum first.")
    else:
        selected_curriculum = st.selectbox(
            "Select Curriculum",
            options=[c['technology'] for c in curricula],
            key="selected_curriculum"
        )
        if selected_curriculum:
            qb_technology = selected_curriculum
            st.write(f"Selected Technology: {qb_technology}")
            num_questions = st.number_input("Number of Questions", min_value=1, value=10, key="num_questions")
            question_type = st.selectbox("Question Type", ["multiple-choice", "subjective", "fill-in-the-blank"], key="question_type")
            difficulty = st.selectbox("Difficulty", ["Easy", "Medium", "Hard"], key="question_difficulty")
        
            if st.button("Generate Question Bank", key="generate_qb_button"):
                curriculum_content = get_curriculum_text(qb_technology)
                if curriculum_content:
                    questions, options, correct_answers = generate_questions(curriculum_content, num_questions, question_type)
        
                    print("Arguments being passed to save_question_bank:")
                    print("qb_technology:", qb_technology)
                    print("topics:", [])
                    print("questions:", '\n'.join(questions))
                    print("difficulty:", difficulty)
                    print("correct_answers:", '\n'.join(correct_answers))
        
                    # Save the generated questions
                    question_bank_id = save_question_bank(qb_technology, [], '\n'.join(questions), difficulty, '\n'.join(correct_answers))
                    if question_bank_id:
                        st.success(f"Question Bank generated successfully! ID: {question_bank_id}")
        
                        # Store the generated questions and options in session state
                        st.session_state.generated_questions = questions
                        st.session_state.generated_options = options
                        st.session_state.generated_qb_id = question_bank_id
                    else:
                        st.error("Failed to save question bank")
                else:
                    st.error("Failed to retrieve curriculum content")
   with tab3:  
      st.subheader("View Generated Questions")  
      question_banks = get_all_question_banks()  
      if not question_banks:  
        st.info("No question banks available yet.")  
      else:  
        selected_qb = st.selectbox(  
           "Select Question Bank",  
           options=[(qb['id'], f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],  
           format_func=lambda x: f"ID: {x[0]} - {x[1]}",  
           key="view_qb_select"  
        )  
  
        if selected_qb:  
           qb_id, _ = selected_qb  
           qb_details = next((qb for qb in question_banks if qb['id'] == qb_id), None)  
           if qb_details:  
              st.write(f"Technology: {qb_details['technology']}")  
              st.write(f"Difficulty: {qb_details['difficulty']}")  
              st.write("Questions:")  
  
              # Check if these are the recently generated questions  
              if 'generated_qb_id' in st.session_state and st.session_state.generated_qb_id == qb_id:  
                for i, (question, options) in enumerate(zip(st.session_state.generated_questions, st.session_state.generated_options), 1):  
                   st.write(f"Q{i}: {question}")  
                   for j, option in enumerate(options):  
                      st.write(f"{chr(65+j)}) {option}")  
                   st.write() # Add a blank line between questions  
              else:  
                questions = qb_details['questions'].split('\n')  
                for i, question in enumerate(questions, 1):  
                   st.write(f"{i}. {question}")  
  
              if st.button("Edit Question Bank"):  
                st.session_state.editing_qb = qb_details  
                st.rerun()  
  
   with tab4:  
      st.subheader("Review Feedback")  
      feedback = review_feedback()  
      if not feedback:  
        st.info("No feedback available yet.")  
      else:  
        feedback_df = pd.DataFrame(feedback)  
        st.dataframe(feedback_df)  
  
        # Sentiment Analysis Summary  
        sentiment_counts = feedback_df['sentiment'].value_counts()  
        st.subheader("Sentiment Analysis Summary")  
        st.bar_chart(sentiment_counts)  
  
        # Question Bank Feedback Summary  
        st.subheader("Question Bank Feedback Summary")  
        qb_feedback = feedback_df.groupby('question_bank_id')['rating'].mean().reset_index()  
        qb_feedback = qb_feedback.merge(pd.DataFrame(question_banks), left_on='question_bank_id', right_on='id')  
        st.dataframe(qb_feedback[['question_bank_id', 'technology', 'difficulty', 'rating']])  
  
   with tab5:  
      st.subheader("Clear History")  
      if st.button("Clear Curriculum Content History"):  
        connection = create_connection()  
        if connection is not None:  
           cursor = connection.cursor()  
           query = "TRUNCATE TABLE curriculum"  
           cursor.execute(query)  
           connection.commit()  
           cursor.close()  
           connection.close()  
           st.success("Curriculum content history cleared successfully!")  
        else:  
           st.error("Failed to connect to database")  
  
      if st.button("Clear Generated Topics History"):  
        connection = create_connection()  
        if connection is not None:  
           cursor = connection.cursor()  
           query = "TRUNCATE TABLE generated_question_files"  
           cursor.execute(query)  
           connection.commit()  
           cursor.close()  
           connection.close()  
           st.success("Generated topics history cleared successfully!")  
        else:  
           st.error("Failed to connect to database")  
  
      st.subheader("Clear Specific Curriculum History")  
      curricula = get_all_curricula()  
      if curricula:  
        selected_curriculum = st.selectbox(  
           "Select Curriculum to Clear",  
           options=[c['technology'] for c in curricula],  
           key="clear_curriculum_select"  
        )  
        if selected_curriculum:  
           if st.button("Clear Curriculum History"):  
              connection = create_connection()  
              if connection is not None:  
                cursor = connection.cursor()  
                query = "DELETE FROM curriculum WHERE technology = %s"  
                cursor.execute(query, (selected_curriculum,))  
                connection.commit()  
                cursor.close()  
                connection.close()  
                st.success(f"Curriculum history for {selected_curriculum} cleared successfully!")  
              else:  
                st.error("Failed to connect to database")  
  
   with tab6:  
    st.subheader("Download Questions")  
    question_banks = get_all_question_banks()  
    if not question_banks:  
      st.info("No question banks available yet.")  
    else:  
      selected_qb = st.selectbox(  
        "Select Question Bank",  
        options=[(qb['id'], f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],  
        format_func=lambda x: f"ID: {x[0]} - {x[1]}",  
        key="download_qb_select"  
      )  
  
      if selected_qb:  
        qb_id, _ = selected_qb  
        qb_details = next((qb for qb in question_banks if qb['id'] == qb_id), None)  
        if qb_details:  
           st.write(f"Technology: {qb_details['technology']}")  
           st.write(f"Difficulty: {qb_details['difficulty']}")  
  
           file_format = st.selectbox("Select File Format", ["docx", "pdf", "pptx", "doc"])  
           if file_format == "pdf":  
              questions = qb_details['questions'].split('\n')  
              pdf = FPDF()  
              pdf.add_page()  
              pdf.set_font("Arial", size=15)  
              for question in questions:  
                pdf.cell(200, 10, txt=question, ln=True, align='L')  
  
              buffer = io.BytesIO()  
              pdf.output(buffer, 'F')  
              buffer.seek(0)  
  
              st.download_button(  
                label='Download PDF',  
                data=buffer.getvalue(),  
                file_name=f'questions_{qb_id}.pdf',  
                mime='application/pdf'  
              )  
           elif file_format == "docx":  
              questions = qb_details['questions'].split('\n')  
              doc = docx.Document()  
              for question in questions:  
                doc.add_paragraph(question)  
              buffer = io.BytesIO()  
              doc.save(buffer)  
              buffer.seek(0)  
  
              st.download_button(  
                label='Download DOCX',  
                data=buffer.getvalue(),  
                file_name=f'questions_{qb_id}.docx',  
                mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document'  
              )  
           elif file_format == "pptx":  
              questions = qb_details['questions'].split('\n')  
              presentation = Presentation()  
              slide_layout = presentation.slide_layouts[6]  
              slide = presentation.slides.add_slide(slide_layout)  
              for question in questions:  
                slide.shapes.add_textbox(10, 10, 300, 100).text = question  
              buffer = io.BytesIO()  
              presentation.save(buffer)  
              buffer.seek(0)  
  
              st.download_button(  
                label='Download PPTX',  
                data=buffer.getvalue(),  
                file_name=f'questions_{qb_id}.pptx',  
                mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'  
              )  
           elif file_format == "doc":  
              questions = qb_details['questions'].split('\n')  
              buffer = io.BytesIO()  
              for question in questions:  
                buffer.write(question.encode('utf-8'))  
                buffer.write(b'\n')  
              buffer.seek(0)  
  
              st.download_button(  
                label='Download DOC',  
                data=buffer.getvalue(),  
                file_name=f'questions_{qb_id}.doc',  
                mime='application/msword'  
              )  
  
   st.subheader("Downloaded Files")  
   downloaded_files = []  
   for file in os.listdir():  
      if file.startswith("questions_") and file.endswith((".pdf", ".docx", ".pptx", ".doc")):  
        downloaded_files.append(file)  
  
   if not downloaded_files:  
      st.info("No downloaded files available.")  
   else:  
      selected_file = st.selectbox(  
        "Select Downloaded File",  
        options=downloaded_files,  
        key="downloaded_file_select"  
      )  
  
      if selected_file:  
        st.write(f"File: {selected_file}")  
  
        if selected_file.endswith(".pdf"):  
           with open(selected_file, "rb") as f:  
              st.download_button(  
                label='Open PDF',  
                data=f.read(),  
                file_name=selected_file,  
                mime='application/pdf'  
              )  
        elif selected_file.endswith(".docx"):  
           with open(selected_file, "rb") as f:  
              st.download_button(  
                label='Open DOCX',  
                data=f.read(),  
                file_name=selected_file,  
                mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document'  
              )  
        elif selected_file.endswith(".pptx"):  
           with open(selected_file, "rb") as f:  
              st.download_button(  
                label='Open PPTX',  
                data=f.read(),  
                file_name=selected_file,  
                mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'  
              )  
        elif selected_file.endswith(".doc"):  
           with open(selected_file, "r") as f:  
              st.write(f.read())  
  
        if st.button("Modify File"):  
           # Add code to modify the file here  
           st.success("File modified successfully!")
  
   
  
   # Additional feature: Curriculum Overview  
   st.sidebar.subheader("Curriculum Overview")  
   curricula = get_all_curricula()  
   if curricula:  
      curriculum_df = pd.DataFrame(curricula)  
      st.sidebar.dataframe(curriculum_df[['technology', 'topics']])  
   else:  
      st.sidebar.info("No curricula available.")
# Download NLTK data
nltk.download('vader_lexicon')

load_dotenv()
os.environ["GOOGLE_API_KEY"] = "AIzaSyAZ11Tinh63Rs1F0yWniCvNG33Q00xag1o"
def save_question_bank(technology, topics, questions, difficulty, correct_answers):  
   connection = create_connection()  
   if connection is None:  
      return False  
  
   try:  
      cursor = connection.cursor()  
      query = """  
      INSERT INTO question_banks (technology, topics, questions, difficulty, correct_answers)  
      VALUES (%s, %s, %s, %s, %s)  
      """  
      cursor.execute(query, (technology, ','.join(topics), questions, difficulty, correct_answers))  
      connection.commit()  
      question_bank_id = cursor.lastrowid  
      return question_bank_id  
   except mysql.connector.Error as err:  
      st.error(f"Database error: {err}")  
      connection.rollback()  
      return False  
   finally:  
      if connection.is_connected():  
        cursor.close()  
        connection.close()
# MySQL connection
def create_connection():
    try:
        connection = mysql.connector.connect(
            host="localhost",
            user="root",
            password="manohar9452",
            database="QB"
        )
        ensure_table_exists(connection)
        return connection
    except Error as e:
        st.error(f"Error connecting to MySQL: {e}")
        return None

genai.configure(api_key=os.getenv('GOOGLE_API_KEY'))
model = genai.GenerativeModel('gemini-pro')

# Utility functions
def extract_text_from_file(file):
    file_extension = os.path.splitext(file.name)[1].lower()
    if file_extension == '.pdf':
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
    elif file_extension == '.docx':
        doc = docx.Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
    elif file_extension == '.txt':
        text = file.getvalue().decode('utf-8')
    else:
        raise ValueError("Unsupported file format")
    return text

def generate_questions(text, num_questions=5, question_type="multiple-choice"):
    if question_type == "multiple-choice":
        prompt = f"Generate {num_questions} multiple-choice questions based on the following text:\n\n{text}\n\nProvide the questions and options in the following format:\n\nQ1: [Question]\nA) [Option 1]\nB) [Option 2]\nC) [Option 3]\nD) [Option 4]\n\nQ2: [Question]\nA) [Option 1]\nB) [Option 2]\nC) [Option 3]\nD) [Option 4]\n\n..."
    elif question_type == "subjective":
        prompt = f"Generate {num_questions} subjective questions based on the following text:\n\n{text}\n\nProvide the questions in the following format:\n\nQ1: [Question]\n\nQ2: [Question]\n\n..."
    elif question_type == "fill-in-the-blank":
        prompt = f"Generate {num_questions} fill-in-the-blank questions based on the following text:\n\n{text}\n\nProvide the questions and correct answers in the following format:\n\nQ1: [Question]\nA: [Correct Answer]\n\nQ2: [Question]\nA: [Correct Answer]\n\n..."
    else:
        raise ValueError("Invalid question type")

    response = model.generate_content(prompt)
    generated_text = response.text

    questions = []
    options = []
    correct_answers = []

    lines = [line.strip() for line in generated_text.split('\n') if line.strip()]

    i = 0
    while i < len(lines):
        if lines[i].startswith('Q'):
            question = lines[i].split(': ', 1)[1]
            questions.append(question)
            if question_type == "multiple-choice":
                options_list = []
                correct_answer = None
                for j in range(i+1, min(i+4, len(lines))):
                    if lines[j].startswith(('A)', 'B)', 'C)', 'D)')):
                        option = lines[j].split(') ', 1)[1]
                        options_list.append(option)
                        if lines[j].startswith('A)'):
                            correct_answer = option
                options.append(options_list)
                correct_answers.append(correct_answer)
                i = j + 1
            elif question_type == "fill-in-the-blank":
                if i+1 < len(lines) and lines[i+1].startswith('A:'):
                    options.append([lines[i+1].split(': ', 1)[1]])
                    correct_answers.append(lines[i+1].split(': ', 1)[1])
                    i += 2
                else:
                    options.append([""])
                    correct_answers.append("")
                    i += 1
            else: # subjective
                options.append([])
                correct_answers.append("")
                i += 1
        else:
            i += 1

    return questions[:num_questions], options[:num_questions], correct_answers[:num_questions]

def ensure_table_exists(connection):
    try:
        cursor = connection.cursor()

        # Check if the table exists
        cursor.execute("SHOW TABLES LIKE 'generated_question_files'")
        result = cursor.fetchone()

        if not result:
            # Table doesn't exist, so create it
            cursor.execute("""
            CREATE TABLE generated_question_files (
                id INT AUTO_INCREMENT PRIMARY KEY,
                technology VARCHAR(255) NOT NULL,
                topics TEXT NOT NULL,
                questions TEXT NOT NULL,
                correct_answers TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
            """)
            connection.commit()
            print("Table 'generated_question_files' created successfully.")
        else:
            print("Table 'generated_question_files' already exists.")

        cursor.close()
    except Error as e:
        print(f"Error ensuring table exists: {e}")

def review_feedback():
    # Simulated feedback data. Replace this with real data from your database or source.
    feedback_data = [
        {"question_bank_id": 1, "sentiment": "positive", "rating": 4.5},
        {"question_bank_id": 1, "sentiment": "negative", "rating": 2.0},
        {"question_bank_id": 2, "sentiment": "neutral", "rating": 3.0},
        {"question_bank_id": 2, "sentiment": "positive", "rating": 4.0}
    ]

    return feedback_data

def analyze_sentiment(text):
    sia = SentimentIntensityAnalyzer()
    sentiment_score = sia.polarity_scores(text)['compound']
    if sentiment_score > 0.05:
        return 'Positive'
    elif sentiment_score < -0.05:
        return 'Negative'
    else:
        return 'Neutral'

def display_questions(questions, options, correct_answers):
    for i, question in enumerate(questions):
        st.write(question)
        if options[i]:
            st.write("Options:")
            for j, option in enumerate(options[i]):
                st.write(f"{chr(65+j)}) {option}")
        st.write(f"Correct Answer: {correct_answers[i]}")
        st.write("")
# User Authentication Functions
def login_user(username, password):
    connection = create_connection()
    if connection is None:
        return None

    cursor = connection.cursor(dictionary=True)
    query = "SELECT * FROM users WHERE username = %s"
    cursor.execute(query, (username,))
    user = cursor.fetchone()
    cursor.close()
    connection.close()

    if user and check_password_hash(user['password'], password):
        return user
    return None

def register_user(email, username, password, role):
    connection = create_connection()
    if connection is None:
        return False

    cursor = connection.cursor()
    check_query = "SELECT * FROM users WHERE username = %s"
    cursor.execute(check_query, (username,))
    existing_user = cursor.fetchone()

    if existing_user:
        cursor.close()
        connection.close()
        return False

    hashed_password = generate_password_hash(password)
    insert_query = "INSERT INTO users (email, username, password, role) VALUES (%s, %s, %s, %s)"
    cursor.execute(insert_query, (email, username, hashed_password, role))
    connection.commit()
    cursor.close()
    connection.close()
    return True

# Administrator Functions
def get_system_stats():
    connection = create_connection()
    if connection is None:
        return None

    cursor = connection.cursor(dictionary=True)
    stats = {}

    tables = ['users', 'question_banks', 'learning_plans', 'feedback']
    for table in tables:
        query = f"SELECT COUNT(*) as count FROM {table}"
        cursor.execute(query)
        result = cursor.fetchone()
        stats[table] = result['count']

    cursor.close()
    connection.close()
    return stats

def get_all_users():
    connection = create_connection()
    if connection is None:
        return []

    cursor = connection.cursor(dictionary=True)
    query = "SELECT username, email, role FROM users"
    cursor.execute(query)
    users = cursor.fetchall()
    cursor.close()
    connection.close()
    return users

def update_user_role(username, new_role):
    connection = create_connection()
    if connection is None:
        return False

    cursor = connection.cursor()
    query = "UPDATE users SET role = %s WHERE username = %s"
    cursor.execute(query, (new_role, username))
    connection.commit()
    affected_rows = cursor.rowcount
    cursor.close()
    connection.close()
    return affected_rows > 0

# Trainer Functions
def upload_curriculum(technology, topics, content):
    connection = create_connection()
    if connection is None:
        return False

    try:
        cursor = connection.cursor()

        # Check if the 'topics' column can accommodate the data
        cursor.execute("SELECT CHARACTER_MAXIMUM_LENGTH FROM INFORMATION_SCHEMA.COLUMNS WHERE TABLE_NAME = 'curriculum' AND COLUMN_NAME = 'topics'")
        result = cursor.fetchone()
        max_topics_length = result[0] if result else 65535 # Default to TEXT max length if not specified

        # Truncate the topics if they are too long
        topics_str = ','.join(topics)
        if len(topics_str) > max_topics_length:
            topics = topics[:int(max_topics_length / len(topics))]
            topics_str = ','.join(topics)
            st.warning(f"The topics were truncated to fit the 'topics' column (max length: {max_topics_length} characters).")

        # Determine the content type (file-like object or string)
        if hasattr(content, 'read'):
            # File-like object
            content_text = content.read().decode('utf-8')
        else:
            # String
            content_text = content

        # Generate questions from the content
        questions, options, correct_answers = generate_questions(content_text)

        # Convert questions, options and correct_answers to strings
        questions_str = '|||'.join(questions)
        options_str = '|||'.join(['###'.join(option) for option in options]) # Using '###' as separator for options
        correct_answers_str = '|||'.join([','.join(map(str, ans)) if isinstance(ans, list) else str(ans) for ans in correct_answers])

        # Insert or update the curriculum
        curriculum_query = """
        INSERT INTO curriculum (technology, topics)
        VALUES (%s, %s)
        ON DUPLICATE KEY UPDATE topics = VALUES(topics)
        """
        cursor.execute(curriculum_query, (technology, topics_str))

        # Insert the generated questions into the generated_question_files table
        questions_query = """
        INSERT INTO generated_question_files (technology, topics, questions, options, correct_answers)
        VALUES (%s, %s, %s, %s, %s)
        """
        cursor.execute(questions_query, (technology, topics_str, questions_str, options_str, correct_answers_str))

        connection.commit()
        return True
    except mysql.connector.Error as err:
        st.error(f"Database error: {err}")
        connection.rollback()
        return False
    finally:
        if connection.is_connected():
            cursor.close()
            connection.close()

def get_curriculum_text(technology):
    connection = create_connection()
    if connection is None:
        return None

    try:
        cursor = connection.cursor(dictionary=True)
        query = "SELECT topics FROM curriculum WHERE technology = %s"
        cursor.execute(query, (technology,))
        result = cursor.fetchone()
        cursor.nextset() # Consume any unread results

        if result:
            return result['topics']
        else:
            st.error(f"No curriculum content found for technology: {technology}")
            return None
    except mysql.connector.Error as err:
        st.error(f"Database error: {err}")
        return None
    finally:
        if cursor:
            cursor.close()
        if connection.is_connected():
            connection.close()

def save_question_bank(technology, topics, questions, difficulty, correct_answers):  
   connection = create_connection()  
   if connection is None:  
      return False  
  
   try:  
      cursor = connection.cursor()  
      query = """  
      INSERT INTO question_banks (technology, topics, questions, difficulty, correct_answers)  
      VALUES (%s, %s, %s, %s, %s)  
      """  
      cursor.execute(query, (technology, ','.join(topics), questions, difficulty, correct_answers))  
      connection.commit()  
      question_bank_id = cursor.lastrowid  
      return question_bank_id  
   except mysql.connector.Error as err:  
      st.error(f"Database error: {err}")  
      connection.rollback()  
      return False  
   finally:  
      if connection.is_connected():  
        cursor.close()  
        connection.close()

def get_topics_for_technology(technology):
    connection = create_connection()
    if connection is None:
        return None

    try:
        cursor = connection.cursor(dictionary=True)
        query = "SELECT topics FROM curriculum WHERE technology = %s"
        cursor.execute(query, (technology,))
        result = cursor.fetchone()
        cursor.nextset() # Consume any unread results

        if result:
            return result['topics'].split(',')
        else:
            st.error(f"No topics found for technology: {technology}")
            return None
    except mysql.connector.Error as err:
        st.error(f"Database error: {err}")
        return None
    finally:
        if cursor:
            cursor.close()
        if connection.is_connected():
            connection.close()

def get_all_curricula():
    connection = create_connection()
    if connection is None:
        return None

    try:
        cursor = connection.cursor(dictionary=True)
        query = "SELECT technology, topics FROM curriculum"
        cursor.execute(query)
        results = cursor.fetchall()
        cursor.nextset() # Consume any unread results

        return results
    except mysql.connector.Error as err:
        st.error(f"Database error: {err}")
        return None
    finally:
        if cursor:
            cursor.close()
        if connection.is_connected():
            connection.close()

def get_all_question_banks():
    connection = create_connection()
    if connection is None:
        return None

    try:
        cursor = connection.cursor(dictionary=True)
        query = "SELECT * FROM question_banks"
        cursor.execute(query)
        results = cursor.fetchall()
        cursor.nextset() # Consume any unread results

        return results
    except mysql.connector.Error as err:
        st.error(f"Database error: {err}")
        return None
    finally:
        if cursor:
            cursor.close()
        if connection.is_connected():
            connection.close()

# Employee Functions
def get_learning_plan(username):
    connection = create_connection()
    if connection is None:
        return None

    cursor = connection.cursor(dictionary=True)
    query = "SELECT * FROM learning_plans WHERE username = %s"
    cursor.execute(query, (username,))
    learning_plan = cursor.fetchone()
    cursor.close()
    connection.close()
    return learning_plan

def submit_feedback(username, question_bank_id, feedback_text, rating):
    sentiment = analyze_sentiment(feedback_text)
    connection = create_connection()
    if connection is None:
        return False

    cursor = connection.cursor()
    query = "INSERT INTO feedback (username, question_bank_id, feedback_text, rating, sentiment) VALUES (%s, %s, %s, %s, %s)"
    cursor.execute(query, (username, question_bank_id, feedback_text, rating, sentiment))
    connection.commit()
    cursor.close()
    connection.close()
    return True

def take_assessment():
    st.subheader("Take Assessment")
    question_banks = get_all_question_banks()
    if not question_banks:
        st.info("No question banks available yet.")
    else:
        selected_qb = st.selectbox(
            "Select Question Bank",
            options=[(qb['id'], f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],
            format_func=lambda x: f"ID: {x[0]} - {x[1]}",
            key="take_assessment_qb_select"
        )

        if selected_qb:
            qb_id, _ = selected_qb
            qb_details = next((qb for qb in question_banks if qb['id'] == qb_id), None)
            if qb_details:
                questions = qb_details['questions'].split('\n')
                options = []
                correct_answers = []
                question_type = None

                # Retrieve question type from question_banks table
                connection = create_connection()
                cursor = connection.cursor()
                query = "SELECT question_type FROM question_banks WHERE id = %s"
                cursor.execute(query, (qb_id,))
                result = cursor.fetchone()
                if result:
                    question_type = result[0]

                # Retrieve options from question_banks table
                query = "SELECT options FROM question_banks WHERE id = %s"
                cursor.execute(query, (qb_id,))
                result = cursor.fetchone()
                if result:
                    options_list = result[0].split('\n')
                    for i in range(0, len(options_list), 4):
                        options.append(options_list[i:i+4])

                # Retrieve correct answers from question_banks table
                query = "SELECT correct_answers FROM question_banks WHERE id = %s"
                cursor.execute(query, (qb_id,))
                result = cursor.fetchone()
                if result:
                    correct_answers = result[0].split('\n')

                score = 0
                answers = []
                for i, question in enumerate(questions):
                    st.write(question)
                    if question_type == "multiple-choice":
                        st.write("Options:")
                        for j, option in enumerate(options[i]):
                            st.write(f"{chr(65+j)}) {option}")
                        answer = st.selectbox("Select an option", [chr(65+j) for j in range(4)], key=f"question_{i}")
                        answers.append(answer)
                    elif question_type == "fill-in-the-blank":
                        answer = st.text_input("Enter your answer", key=f"question_{i}")
                        answers.append(answer)
                    elif question_type == "subjective":
                        answer = st.text_area("Enter your answer", key=f"question_{i}")
                        answers.append(answer)

                if st.button("Submit"):
                    for i, answer in enumerate(answers):
                        if question_type == "multiple-choice":
                            if answer == correct_answers[i][0]:
                                score += 1
                        elif question_type == "fill-in-the-blank":
                            if answer == correct_answers[i]:
                                score += 1
                        elif question_type == "subjective":
                            if answer == correct_answers[i]:
                                score += 1

                    st.write(f"Your score is {score} out of {len(questions)}")

                    # Save the assessment result
                    connection = create_connection()
                    cursor = connection.cursor()
                    query = "INSERT INTO assessments (username, question_bank_id, score) VALUES (%s, %s, %s)"
                    cursor.execute(query, (st.session_state.username, qb_id, score))
                    connection.commit()
                    connection.close()

def get_available_question_banks(username):
    connection = create_connection()
    if connection is None:
        return []

    cursor = connection.cursor(dictionary=True)
    query = """
    SELECT qb.id, qb.technology, qb.topics
    FROM question_banks qb
    JOIN learning_plans lp ON qb.technology = lp.technology
    WHERE lp.username = %s AND qb.id NOT IN (
        SELECT question_bank_id FROM assessments WHERE username = %s
    )
    """
    cursor.execute(query, (username, username))
    question_banks = cursor.fetchall()
    cursor.close()
    connection.close()
    return question_banks

def get_completed_assessments(username):
    connection = create_connection()
    if connection is None:
        return []

    cursor = connection.cursor(dictionary=True)
    query = """
    SELECT a.id, qb.technology
    FROM assessments a
    JOIN question_banks qb ON a.question_bank_id = qb.id
    WHERE a.username = %s
    """
    cursor.execute(query, (username,))
    completed_assessments = cursor.fetchall()
    cursor.close()
    connection.close()
    return completed_assessments

# Dashboard Functions

def admin_dashboard():  
   st.header("Administrator Dashboard")  
   tab1, tab2, tab3 = st.tabs(["System Stats", "User Management", "Reports"])  
  
   with tab1:  
      st.subheader("System Statistics")  
      stats = get_system_stats()  
      if stats:  
        st.write(f"Total Users: {stats['users']}")  
        st.write(f"Total Question Banks: {stats['question_banks']}")  
        st.write(f"Total Learning Plans: {stats['learning_plans']}")  
        st.write(f"Total Feedback Entries: {stats['feedback']}")  
      else:  
        st.error("Failed to retrieve system statistics")  
  
   with tab2:  
      st.subheader("User Management")  
      users = get_all_users()  
      for user in users:  
        st.write(f"Username: {user['username']}, Email: {user['email']}, Role: {user['role']}")  
        new_role = st.selectbox(f"New Role for {user['username']}", ["Administrator", "Trainer", "Employee"], key=user['username'])  
        if st.button(f"Update Role for {user['username']}"):  
           if update_user_role(user['username'], new_role):  
              st.success(f"Role updated for {user['username']}")  
           else:  
              st.error("Failed to update role")  
  
   with tab3:  
      st.subheader("Generate Reports")  
      report_type = st.selectbox("Select Report Type", ["User Activity", "Question Bank Usage", "Feedback Summary", "Sentiment Analysis"])  
      if st.button("Generate Report"):  
        if report_type == "User Activity":  
           user_activity_report()  
        elif report_type == "Question Bank Usage":  
           question_bank_usage_report()  
        elif report_type == "Feedback Summary":  
           feedback_summary_report()  
        elif report_type == "Sentiment Analysis":  
           sentiment_analysis_report() 
    
   notifications = get_notifications("admin")  
   if notifications:  
      st.sidebar.write("Notifications:")  
      st.sidebar.write(len(notifications))  
      if st.sidebar.button("View Notifications"):  
        st.sidebar.write("Notifications:")  
        for notification in notifications:  
           st.sidebar.write(notification['message'])  
   else:  
      st.sidebar.write("No notifications available.")
def user_activity_report():  
   connection = create_connection()  
   if connection is None:  
      st.error("Failed to connect to database")  
      return  
  
   cursor = connection.cursor(dictionary=True)  
   query = "SELECT username, COUNT(*) as num_assessments FROM assessments GROUP BY username"  
   cursor.execute(query)  
   results = cursor.fetchall()  
   cursor.close()  
   connection.close()  
  
   if results:  
      user_activity = {}  
      user_activity_table = []  
      for result in results:  
        username = result['username']  
        num_assessments = result['num_assessments']  
        user_activity[username] = num_assessments  
        user_activity_table.append({  
           'Username': username,  
           'Number of Assessments': num_assessments  
        })  
  
      st.write("User Activity Report:")  
      st.write("User Activity Table:")  
      st.table(pd.DataFrame(user_activity_table))  
      st.write("User Activity Chart:")  
      st.bar_chart(user_activity)  
   else:  
      st.error("No user activity data available")  
  
def question_bank_usage_report():  
   connection = create_connection()  
   if connection is None:  
      st.error("Failed to connect to database")  
      return  
  
   cursor = connection.cursor(dictionary=True)  
   query = "SELECT question_bank_id, COUNT(*) as num_assessments FROM assessments GROUP BY question_bank_id"  
   cursor.execute(query)  
   results = cursor.fetchall()  
   cursor.close()  
   connection.close()  
  
   if results:  
      question_bank_usage = {}  
      question_bank_usage_table = []  
      for result in results:  
        question_bank_id = result['question_bank_id']  
        num_assessments = result['num_assessments']  
        question_bank_usage[question_bank_id] = num_assessments  
        question_bank_usage_table.append({  
           'Question Bank ID': question_bank_id,  
           'Number of Assessments': num_assessments  
        })  
  
      st.write("Question Bank Usage Report:")  
      st.write("Question Bank Usage Table:")  
      st.table(pd.DataFrame(question_bank_usage_table))  
      st.write("Question Bank Usage Chart:")  
      st.bar_chart(question_bank_usage)  
   else:  
      st.error("No question bank usage data available")  
  
def sentiment_analysis_report():  
   connection = create_connection()  
   if connection is None:  
      st.error("Failed to connect to database")  
      return  
  
   cursor = connection.cursor(dictionary=True)  
   query = "SELECT sentiment, COUNT(*) as num_feedback FROM feedback GROUP BY sentiment"  
   cursor.execute(query)  
   results = cursor.fetchall()  
   cursor.close()  
   connection.close()  
  
   if results:  
      sentiment_analysis = {}  
      sentiment_analysis_table = []  
      for result in results:  
        sentiment = result['sentiment']  
        num_feedback = result['num_feedback']  
        sentiment_analysis[sentiment] = num_feedback  
        sentiment_analysis_table.append({  
           'Sentiment': sentiment,  
           'Number of Feedback': num_feedback  
        })  
  
      st.write("Sentiment Analysis Report:")  
      st.write("Sentiment Analysis Table:")  
      st.table(pd.DataFrame(sentiment_analysis_table))  
      st.write("Sentiment Analysis Chart:")  
      st.bar_chart(sentiment_analysis)  
   else:  
      st.error("No sentiment analysis data available")  
  
def feedback_summary_report():  
   connection = create_connection()  
   if connection is None:  
      st.error("Failed to connect to database")  
      return  
  
   cursor = connection.cursor(dictionary=True)  
   query = "SELECT question_bank_id, AVG(rating) as avg_rating, COUNT(*) as num_feedback FROM feedback GROUP BY question_bank_id"  
   cursor.execute(query)  
   results = cursor.fetchall()  
   cursor.close()  
   connection.close()  
  
   if results:  
      feedback_summary = {}  
      feedback_summary_table = []  
      for result in results:  
        question_bank_id = result['question_bank_id']  
        avg_rating = result['avg_rating']  
        num_feedback = result['num_feedback']  
        feedback_summary[question_bank_id] = avg_rating  
        feedback_summary_table.append({  
           'Question Bank ID': question_bank_id,  
           'Average Rating': avg_rating,  
           'Number of Feedback': num_feedback  
        })  
  
      st.write("Feedback Summary Report:")  
      st.write("Feedback Summary Table:")  
      st.table(pd.DataFrame(feedback_summary_table))  
      st.write("Feedback Summary Chart:")  
      st.bar_chart(feedback_summary)  
   else:  
      st.error("No feedback data available")

def feedback_received(feedback):  
   connection = create_connection()  
   if connection is None:  
      return False  
  
   try:  
      cursor = connection.cursor()  
      query = "INSERT INTO feedback (feedback) VALUES (%s)"  
      cursor.execute(query, (feedback,))  
      connection.commit()  
      cursor.close()  
      connection.close()  
  
      print("Feedback received!")  
  
      # Send notification to admin  
      recipient_role = "admin"  
      message = f"New feedback received: {feedback}"  
      print(f"Sending notification to {recipient_role}: {message}")  
      send_notification(recipient_role, message)  
  
      # Send notification to trainer  
      recipient_role = "trainer"  
      message = f"New feedback received: {feedback}"  
      print(f"Sending notification to {recipient_role}: {message}")  
      send_notification(recipient_role, message)  
  
      # Send notification to employee  
      recipient_role = "employee"  
      message = f"New feedback received: {feedback}"  
      print(f"Sending notification to {recipient_role}: {message}")  
      send_notification(recipient_role, message)  
  
      return True  
   except mysql.connector.Error as err:  
      st.error(f"Database error: {err}")  
      connection.rollback()  
      return False

def trainer_dashboard():  
   st.header("Trainer Dashboard")  
   tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs(["Upload Curriculum", "Generate Question Bank", "View Questions", "Review Feedback", "Clear History", "Download Questions"])  
  
   with tab1:  
      st.subheader("Upload Curriculum")  
      technology = st.text_input("Technology", key="upload_technology")  
      topics = st.text_area("Topics (one per line)", key="upload_topics")  
      uploaded_file = st.file_uploader("Upload curriculum file", type=['txt', 'pdf', 'docx'], key="curriculum_file")  
  
      if st.button("Upload Curriculum", key="upload_curriculum_button"):  
        topic_list = [topic.strip() for topic in topics.split('\n') if topic.strip()]  
        content = ""  
  
        if uploaded_file is not None:  
           try:  
              file_content = extract_text_from_file(uploaded_file)  
              content = file_content  
              topic_list.extend([topic.strip() for topic in file_content.split('\n') if topic.strip()])  
              topic_list = list(set(topic_list)) # Remove duplicates  
           except ValueError as e:  
              st.error(f"Error processing file: {str(e)}")  
              return  
  
        if upload_curriculum(technology, topic_list, content):  
           st.success("Curriculum uploaded successfully!")  
        else:  
           st.error("Failed to upload curriculum")  
  
   with tab2:  
      st.subheader("Generate Question Bank")  
  
      curricula = get_all_curricula()  
  
      if not curricula:  
        st.warning("No curricula available. Please upload a curriculum first.")  
      else:  
        selected_curriculum = st.selectbox(  
           "Select Curriculum",  
           options=[c['technology'] for c in curricula],  
           key="selected_curriculum"  
        )  
        if selected_curriculum:  
           qb_technology = selected_curriculum  
           st.write(f"Selected Technology: {qb_technology}")  
           num_questions = st.number_input("Number of Questions", min_value=1, value=10, key="num_questions")  
           question_type = st.selectbox("Question Type", ["multiple-choice", "subjective", "fill-in-the-blank"], key="question_type")  
           difficulty = st.selectbox("Difficulty", ["Easy", "Medium", "Hard"], key="question_difficulty")  
  
           if st.button("Generate Question Bank", key="generate_qb_button"):
            curriculum_content = get_curriculum_text(qb_technology)
            if curriculum_content:
               questions, options, correct_answers = generate_questions(curriculum_content, num_questions, question_type)
            
               print("Arguments being passed to save_question_bank:")
               print("qb_technology:", qb_technology)
               print("topics:", [])
               print("questions:", '\n'.join(questions))
               print("difficulty:", difficulty)
               print("correct_answers:", '\n'.join(correct_answers))
            
               # Save the generated questions
               question_bank_id = save_question_bank(qb_technology, [], '\n'.join(questions), difficulty, '\n'.join(correct_answers))
               if question_bank_id:
                     st.success(f"Question Bank generated successfully! ID: {question_bank_id}")
            
                     # Store the generated questions and options in session state
                     st.session_state.generated_questions = questions
                     st.session_state.generated_options = options
                     st.session_state.generated_qb_id = question_bank_id
               else:
                     st.error("Failed to save question bank")
            else:
               st.error("Failed to retrieve curriculum content") 
  
   with tab3:  
      st.subheader("View Generated Questions")  
      question_banks = get_all_question_banks()  
      if not question_banks:  
        st.info("No question banks available yet.")  
      else:  
        selected_qb = st.selectbox(  
           "Select Question Bank",  
           options=[(qb['id'], f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],  
           format_func=lambda x: f"ID: {x[0]} - {x[1]}",  
           key="view_qb_select"  
        )  
  
        if selected_qb:  
           qb_id, _ = selected_qb  
           qb_details = next((qb for qb in question_banks if qb['id'] == qb_id), None)  
           if qb_details:  
              st.write(f"Technology: {qb_details['technology']}")  
              st.write(f"Difficulty: {qb_details['difficulty']}")  
              st.write("Questions:")  
  
              # Check if these are the recently generated questions  
              if 'generated_qb_id' in st.session_state and st.session_state.generated_qb_id == qb_id:  
                for i, (question, options) in enumerate(zip(st.session_state.generated_questions, st.session_state.generated_options), 1):  
                   st.write(f"Q{i}: {question}")  
                   for j, option in enumerate(options):  
                      st.write(f"{chr(65+j)}) {option}")  
                   st.write() # Add a blank line between questions  
              else:  
                questions = qb_details['questions'].split('\n')  
                for i, question in enumerate(questions, 1):  
                   st.write(f"{i}. {question}")  
  
              if st.button("Edit Question Bank"):  
                st.session_state.editing_qb = qb_details  
                st.rerun()  
  
   with tab4:  
      st.subheader("Review Feedback")  
      feedback = review_feedback()  
      if not feedback:  
        st.info("No feedback available yet.")  
      else:  
        feedback_df = pd.DataFrame(feedback)  
        st.dataframe(feedback_df)  
  
        # Sentiment Analysis Summary  
        sentiment_counts = feedback_df['sentiment'].value_counts()  
        st.subheader("Sentiment Analysis Summary")  
        st.bar_chart(sentiment_counts)  
  
        # Question Bank Feedback Summary  
        st.subheader("Question Bank Feedback Summary")  
        qb_feedback = feedback_df.groupby('question_bank_id')['rating'].mean().reset_index()  
        qb_feedback = qb_feedback.merge(pd.DataFrame(question_banks), left_on='question_bank_id', right_on='id')  
        st.dataframe(qb_feedback[['question_bank_id', 'technology', 'difficulty', 'rating']])  
  
   with tab5:  
      st.subheader("Clear History")  
      if st.button("Clear Curriculum Content History"):  
        connection = create_connection()  
        if connection is not None:  
           cursor = connection.cursor()  
           query = "TRUNCATE TABLE curriculum"  
           cursor.execute(query)  
           connection.commit()  
           cursor.close()  
           connection.close()  
           st.success("Curriculum content history cleared successfully!")  
        else:  
           st.error("Failed to connect to database")  
  
      if st.button("Clear Generated Topics History"):  
        connection = create_connection()  
        if connection is not None:  
           cursor = connection.cursor()  
           query = "TRUNCATE TABLE generated_question_files"  
           cursor.execute(query)  
           connection.commit()  
           cursor.close()  
           connection.close()  
           st.success("Generated topics history cleared successfully!")  
        else:  
           st.error("Failed to connect to database")  
  
   with tab6:   
    st.subheader("Download Questions")   
    question_banks = get_all_question_banks()   
    if not question_banks:   
        st.info("No question banks available yet.")   
    else:   
        selected_qb = st.selectbox(   
        "Select Question Bank",   
        options=[(qb['id'], f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],   
        format_func=lambda x: f"ID: {x[0]} - {x[1]}",   
        key="download_qb_select"   
        )   
      
        if selected_qb:   
          qb_id, _ = selected_qb   
          qb_details = next((qb for qb in question_banks if qb['id'] == qb_id), None)   
          if qb_details:   
            st.write(f"Technology: {qb_details['technology']}")   
            st.write(f"Difficulty: {qb_details['difficulty']}")   
        
            st.subheader("Download Documents")   
            doc_type = st.selectbox("Select Document Type", ["Question Bank", "Personalized Learning Plan", "Feedback Report", "Curriculum Mapping Document", "Issue Resolution Report", "Assessment Completion Report"])   
        
            if doc_type == "Question Bank":   
              doc = generate_question_bank_document(qb_id)   
              if doc:   
                st.download_button("Download Question Bank", doc, file_name=f"question_bank_{qb_id}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")   
              else:   
                st.error("Failed to generate question bank document")   
        
            elif doc_type == "Personalized Learning Plan":   
              st.write("You can check your learning plan in your Employee Dashboard.")  
        
            elif doc_type == "Feedback Report":   
              doc = generate_feedback_report_document(qb_id)   
              if doc:   
                st.download_button("Download Feedback Report", doc, file_name=f"feedback_report_{qb_id}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")   
              else:   
                st.error("Failed to generate feedback report document")   
        
            elif doc_type == "Curriculum Mapping Document":   
              technology = st.text_input("Enter Technology")   
              doc = generate_curriculum_mapping_document(technology)   
              if doc:   
                st.download_button("Download Curriculum Mapping Document", doc, file_name=f"curriculum_mapping_{technology}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")   
              else:   
                st.error("Failed to generate curriculum mapping document")   
        
            elif doc_type == "Issue Resolution Report":   
              doc = generate_issue_resolution_report_document(qb_id)   
              if doc:   
                st.download_button("Download Issue Resolution Report", doc, file_name=f"issue_resolution_report_{qb_id}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")   
              else:   
                st.error("Failed to generate issue resolution report document")   
        
            elif doc_type == "Assessment Completion Report":   
              username = st.text_input("Enter Username")   
              doc = generate_assessment_completion_report_document(username)   
              if doc:   
                st.download_button("Download Assessment Completion Report", doc, file_name=f"assessment_completion_report_{username}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")   
              else:   
                st.error("Failed to generate assessment completion report document")
    
    # Additional feature: Curriculum Overview
    st.sidebar.subheader("Curriculum Overview")
    curricula = get_all_curricula()
    if curricula:
        curriculum_df = pd.DataFrame(curricula)
        st.sidebar.dataframe(curriculum_df[['technology', 'topics']])
    else:
        st.sidebar.info("No curricula available.")

def get_curriculum_history():  
  connection = create_connection()  
  if connection is None:  
    return None  
  
  cursor = connection.cursor(dictionary=True)  
  query = "SELECT id, technology, topics FROM curriculum"  
  cursor.execute(query)  
  results = cursor.fetchall()  
  cursor.close()  
  connection.close()  
  
  return results

def generate_question_bank_document(qb_id):  
   connection = create_connection()  
   if connection is None:  
      return None  
  
   cursor = connection.cursor(dictionary=True)  
   query = "SELECT questions FROM question_banks WHERE id = %s"  
   cursor.execute(query, (qb_id,))  
   result = cursor.fetchone()  
   cursor.close()  
   connection.close()  
  
   if result:  
      questions = result['questions'].split('\n')  
      doc = docx.Document()  
      for question in questions:  
        doc.add_paragraph(question)  
      buffer = io.BytesIO()  
      doc.save(buffer)  
      buffer.seek(0)  
      return buffer.getvalue()  
   else:  
      return None  
  
def generate_personalized_learning_plan_document(username):   
  connection = create_connection()   
  if connection is None:   
    return None   
   
  cursor = connection.cursor(dictionary=True)   
  query = "SELECT id FROM question_banks WHERE id IN (SELECT question_bank_id FROM learning_plans WHERE username = %s)"   
  cursor.execute(query, (username,))   
  result = cursor.fetchone()   
  cursor.close()   
  connection.close()   
   
  if result:   
    qb_id = result['id']   
    learning_plan = prepare_learning_plan(qb_id, username)   
    if learning_plan:   
      doc = docx.Document()   
      doc.add_paragraph(f"Technology: {learning_plan['technology']}")   
      doc.add_paragraph(f"Start Date: {learning_plan['start_date']}")   
      doc.add_paragraph(f"End Date: {learning_plan['end_date']}")   
      doc.add_paragraph(f"Status: {learning_plan['status']}")   
      doc.add_paragraph(f"Estimated Time: {learning_plan['estimated_time']} hours")   
   
      buffer = io.BytesIO()   
      doc.save(buffer)   
      buffer.seek(0)   
      return buffer.getvalue()   
    else:   
      return None   
  else:   
    return None
  
def generate_feedback_report_document(qb_id):  
   connection = create_connection()  
   if connection is None:  
      return None  
  
   cursor = connection.cursor(dictionary=True)  
   query = "SELECT feedback_text, rating FROM feedback WHERE question_bank_id = %s"  
   cursor.execute(query, (qb_id,))  
   results = cursor.fetchall()  
   cursor.close()  
   connection.close()  
  
   if results:  
      doc = docx.Document()  
      for result in results:  
        doc.add_paragraph(result['feedback_text'])  
        doc.add_paragraph(str(result['rating']))  
      buffer = io.BytesIO()  
      doc.save(buffer)  
      buffer.seek(0)  
      return buffer.getvalue()  
   else:  
      return None  
  
def generate_curriculum_mapping_document(technology):  
   connection = create_connection()  
   if connection is None:  
      return None  
  
   cursor = connection.cursor(dictionary=True)  
   query = "SELECT topics FROM curriculum WHERE technology = %s"  
   cursor.execute(query, (technology,))  
   result = cursor.fetchone()  
   cursor.nextset()  # Consume any unread results  
   cursor.close()  
   connection.close()  
  
   if result:  
      topics = result['topics'].split(',')  
      doc = docx.Document()  
      for topic in topics:  
        doc.add_paragraph(topic)  
      buffer = io.BytesIO()  
      doc.save(buffer)  
      buffer.seek(0)  
      return buffer.getvalue()  
   else:  
      return None
  
def generate_issue_resolution_report_document(qb_id):  
   connection = create_connection()  
   if connection is None:  
      return None  
  
   cursor = connection.cursor(dictionary=True)  
   query = "SELECT issue, resolution FROM issue_resolution WHERE question_bank_id = %s"  
   cursor.execute(query, (qb_id,))  
   results = cursor.fetchall()  
   cursor.close()  
   connection.close()  
  
   if results:  
      doc = docx.Document()  
      for result in results:  
        doc.add_paragraph(result['issue'])  
        doc.add_paragraph(result['resolution'])  
      buffer = io.BytesIO()  
      doc.save(buffer)  
      buffer.seek(0)  
      return buffer.getvalue()  
   else:  
      return None  
  
def generate_assessment_completion_report_document(username):  
   connection = create_connection()  
   if connection is None:  
      return None  
  
   cursor = connection.cursor(dictionary=True)  
   query = "SELECT score, completed_at FROM assessments WHERE username = %s"  
   cursor.execute(query, (username,))  
   results = cursor.fetchall()  
   cursor.close()  
   connection.close()  
  
   if results:  
      doc = docx.Document()  
      for result in results:  
        doc.add_paragraph(str(result['score']))  
        doc.add_paragraph(str(result['completed_at']))  
      buffer = io.BytesIO()  
      doc.save(buffer)  
      buffer.seek(0)  
      return buffer.getvalue()  
   else:  
      return None
def send_notification(recipient_role, message):  
   connection = create_connection()  
   if connection is None:  
      return False  
  
   try:  
      cursor = connection.cursor()  
      query = "INSERT INTO notifications (recipient_role, message) VALUES (%s, %s)"  
      cursor.execute(query, (recipient_role, message))  
      connection.commit()  
      cursor.close()  
      connection.close()  
      return True  
   except mysql.connector.Error as err:  
      st.error(f"Database error: {err}")  
      connection.rollback()  
      return False  
  
def get_notifications(recipient_role):  
   connection = create_connection()  
   if connection is None:  
      return []  
  
   try:  
      cursor = connection.cursor(dictionary=True)  
      query = "SELECT message FROM notifications WHERE recipient_role = %s ORDER BY created_at DESC"  
      cursor.execute(query, (recipient_role,))  
      notifications = cursor.fetchall()  
      cursor.close()  
      connection.close()  
      return notifications  
   except mysql.connector.Error as err:  
      st.error(f"Database error: {err}")  
      return [] 

# New function to update question bank
def update_question_bank(qb_id, new_questions):
    connection = create_connection()
    if connection is None:
        return False

    try:
        cursor = connection.cursor()
        query = "UPDATE question_banks SET questions = %s WHERE id = %s"
        cursor.execute(query, (new_questions, qb_id))
        connection.commit()
        return True
    except mysql.connector.Error as err:
        st.error(f"Database error: {err}")
        return False
    finally:
        if connection.is_connected():
            cursor.close()
            connection.close()

def update_learning_plan_status(qb_id, username, new_status):   
  connection = create_connection()   
  if connection is None:   
    return False   
   
  cursor = connection.cursor(dictionary=True)   
  query = "SELECT technology FROM question_banks WHERE id = %s"   
  cursor.execute(query, (qb_id,))   
  result = cursor.fetchone()   
  cursor.close()   
  connection.close()   
   
  if result:   
    technology = result['technology']   
    connection = create_connection()   
    if connection is None:   
      return False   
   
    cursor = connection.cursor()   
    query = "SHOW COLUMNS FROM learning_plans LIKE 'status'"   
    cursor.execute(query)   
    result = cursor.fetchone()   
    if not result:   
      query = "ALTER TABLE learning_plans ADD COLUMN status VARCHAR(255)"   
      cursor.execute(query)   
      connection.commit()   
   
    query = "UPDATE learning_plans SET status = %s WHERE username = %s AND technology = %s"   
    cursor.execute(query, (new_status, username, technology))   
    connection.commit()   
    cursor.close()   
    connection.close()   
    return True   
  else:   
    return False

  
def employee_dashboard(username):    
  st.header(f"Welcome, {username}!")    
  tab1, tab2, tab3, tab4 = st.tabs(["Learning Plan", "Prepare from Generated Questions", "Take Assessment", "View Progress"])    
    
  with tab1:    
    st.subheader("Your Learning Plan")    
    question_banks = get_all_question_banks()    
    if question_banks:    
      selected_qb = st.selectbox(    
        "Select Question Bank for Learning Plan",    
        options=[(qb['id'], f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],    
        format_func=lambda x: f"ID: {x[0]} - {x[1]}",    
        key="learning_plan_qb_select"    
      )    
      if selected_qb:    
        qb_id, _ = selected_qb    
        learning_plan = prepare_learning_plan(qb_id, username)    
        if learning_plan:    
          st.write(f"Technology: {learning_plan['technology']}")    
          st.write(f"Start Date: {learning_plan['start_date']}")    
          st.write(f"End Date: {learning_plan['end_date']}")    
          st.write(f"Status: {learning_plan['status']}")    
          st.write(f"Estimated Time to Complete: {learning_plan['estimated_time']} hours")    
    
          # Update status feature    
          new_status = st.selectbox("Update Status", ["In Progress", "Completed"])    
          if st.button("Update Status"):    
           if update_learning_plan_status(qb_id, username, new_status):    
             st.success("Status updated successfully!")    
           else:    
             st.error("Failed to update status")    
    
          # Generate next learning plan feature    
          if st.button("Generate Next Learning Plan"):    
           next_qb_id = get_next_question_bank_id(qb_id)    
           if next_qb_id:    
             next_learning_plan = prepare_learning_plan(next_qb_id, username)    
             if next_learning_plan:    
               st.write(f"Next Technology: {next_learning_plan['technology']}")    
               st.write(f"Next Start Date: {next_learning_plan['start_date']}")    
               st.write(f"Next End Date: {next_learning_plan['end_date']}")    
               st.write(f"Next Status: {next_learning_plan['status']}")    
               st.write(f"Next Estimated Time to Complete: {next_learning_plan['estimated_time']} hours")    
             else:    
               st.error("Failed to generate next learning plan")    
           else:    
             st.error("No next question bank available")    
        else:    
          st.warning("No learning plan assigned yet.")    
    else:    
      st.info("No question banks available yet.")   
    
  with tab2:    
   st.subheader("Prepare from Generated Questions")    
   question_banks = get_all_question_banks()    
   if question_banks:    
    selected_qb = st.selectbox(    
      "Select Question Bank",    
      options=[(qb['id'], f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],    
      format_func=lambda x: f"ID: {x[0]} - {x[1]}",    
      key="prepare_qb_select"    
    )    
    if selected_qb:    
      qb_id, _ = selected_qb    
      qb_details = next((qb for qb in question_banks if qb['id'] == qb_id), None)    
      if qb_details:    
       st.write(f"Technology: {qb_details['technology']}")    
       st.write(f"Difficulty: {qb_details['difficulty']}")    
       st.write("Questions:")    
       questions = qb_details['questions'].split('\n')    
    
       # Translation feature    
       translate_language = st.selectbox("Select Language", [    
        "English",    
        "Spanish",    
        "French",    
        "German",    
        "Chinese",    
        "Japanese",    
        "Korean",    
        "Hindi",    
        "Tamil",    
        "Telugu",    
        "Bengali",    
        "Marathi",    
        "Gujarati",    
        "Punjabi",    
        "Kannada",    
        "Malayalam"    
       ])    
       if translate_language != "English":    
        # Use Google Translator API to translate questions    
        translator = googletrans.Translator()    
        translated_questions = []    
        for question in questions:    
         try:    
          translated_question = translator.translate(question, dest=translate_language).text    
          translated_questions.append(translated_question)    
         except Exception as e:    
          st.error(f"Error translating question: {e}")    
        st.write("Translated Questions:")    
        for i, question in enumerate(translated_questions, 1):    
         st.write(f"{i}. {question}")    
       else:    
        for i, question in enumerate(questions, 1):    
         st.write(f"{i}. {question}")    
      else:    
       st.error("Failed to retrieve question bank details")    
   else:    
    st.info("No question banks available yet.")    
    
  with tab3:  
   st.subheader("Take Assessment")  
   question_banks = get_all_question_banks()  
   if not question_banks:  
      st.info("No question banks available yet.")  
   else:  
      selected_qb = st.selectbox(  
        "Select Question Bank",  
        options=[(qb['id'], f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],  
        format_func=lambda x: f"ID: {x[0]} - {x[1]}",  
        key="take_assessment_qb_select"  
      )  
  
      if selected_qb:  
        qb_id, _ = selected_qb  
        qb_details = next((qb for qb in question_banks if qb['id'] == qb_id), None)  
        if qb_details:  
           st.write(f"Technology: {qb_details['technology']}")  
           st.write(f"Difficulty: {qb_details['difficulty']}")  
  
           questions = qb_details['questions'].split('\n')  
           correct_answers = qb_details['correct_answers'].split('\n')  
  
           user_answers = []  
           for i, question in enumerate(questions, 1):  
              user_answer = st.text_input(f"Q{i}: {question}", key=f"user_answer_{i}")  
              user_answers.append(user_answer)  
  
           if st.button("Submit"):  
              score = 0  
              for i, (user_answer, correct_answer) in enumerate(zip(user_answers, correct_answers)):  
                if user_answer == correct_answer:  
                   st.success(f"Q{i+1}: Correct!")  
                   score += 1  
                else:  
                   st.error(f"Q{i+1}: Incorrect. Correct answer: {correct_answer}")  
              st.write(f"Score: {score}/{len(questions)}")  
    
  with tab4:    
   st.subheader("Your Progress")    
   completed_assessments = get_completed_assessments(username)    
   if completed_assessments:    
    for assessment in completed_assessments:    
      st.write(f"Assessment ID: {assessment['id']}")    
      st.write(f"Technology: {assessment['technology']}")    
      st.write("---")    
   else:    
    st.info("You haven't completed any assessments yet.")    
    
   # Feedback submission    
   st.subheader("Submit Feedback")    
   feedback_text = st.text_area("Your Feedback")    
   rating = st.slider("Rating", 1, 5, 3)    
   if st.button("Submit Feedback"):    
    if submit_feedback(username, selected_qb[0], feedback_text, rating):    
      st.success("Feedback submitted successfully!")    
    else:    
      st.error("Failed to submit feedback")
   notifications = get_notifications("employee")  
   if notifications:  
      st.sidebar.write("Notifications:")  
      for notification in notifications:  
        st.sidebar.write(notification['message'])

def get_generated_questions():
    connection = create_connection()
    if connection is None:
        return None

    try:
        cursor = connection.cursor()
        query = """
        SELECT questions, options, correct_answers
        FROM generated_question_files
        """
        cursor.execute(query)
        result = cursor.fetchone()
        if result:
            return {
                'questions': result[0],
                'options': result[1],
                'correct_answers': result[2]
            }
        else:
            return None
    except mysql.connector.Error as err:
        st.error(f"Database error: {err}")
        return None
    finally:
        if connection.is_connected():
            cursor.close()
            connection.close()

def get_next_question_bank_id(qb_id):   
  connection = create_connection()   
  if connection is None:   
    return None   
   
  cursor = connection.cursor(dictionary=True)   
  query = "SELECT id FROM question_banks WHERE id > %s ORDER BY id ASC LIMIT 1"   
  cursor.execute(query, (qb_id,))   
  result = cursor.fetchone()   
  cursor.close()   
  connection.close()   
  if result:   
    return result['id']   
  else:   
    return None

def prepare_learning_plan(qb_id, username):   
  connection = create_connection()   
  if connection is None:   
    return None   
   
  cursor = connection.cursor(dictionary=True)   
  query = "SELECT questions, difficulty FROM question_banks WHERE id = %s"   
  cursor.execute(query, (qb_id,))   
  result = cursor.fetchone()   
  cursor.close()   
  connection.close()   
   
  if result:   
    questions = result['questions'].split('\n')   
    difficulty = result['difficulty']   
    num_questions = len(questions)   
    estimated_time = calculate_estimated_time(num_questions, difficulty)   
    topics = []   
    for question in questions:   
      topic = question.split(':')[0].strip()   
      if topic not in topics:   
        topics.append(topic)   
   
    # Get the date when the employee updated the status as 'Completed'   
    completed_date = get_completed_date(username)   
   
    # Set the start date to the day after the completed date   
    start_date = (completed_date + timedelta(days=1)).strftime('%Y-%m-%d')   
   
    # Calculate the estimated end date based on the topics length and other analysis   
    estimated_end_date = calculate_estimated_end_date(topics, estimated_time, start_date)   
   
    learning_plan = {   
      'technology': topics,   
      'start_date': start_date,   
      'end_date': estimated_end_date,   
      'status': 'In Progress',   
      'estimated_time': estimated_time   
    }   
   
    return learning_plan   
  else:   
    return None

def get_correct_answers(qb_id):  
   connection = create_connection()  
   if connection is None:  
      return None  
  
   try:  
      cursor = connection.cursor(dictionary=True)  
      query = "SELECT correct_answers FROM question_banks WHERE id = %s"  
      cursor.execute(query, (qb_id,))  
      result = cursor.fetchone()  
      cursor.close()  
      connection.close()  
  
      if result:  
        return result['correct_answers'].split('\n')  
      else:  
        return None  
   except mysql.connector.Error as err:  
      st.error(f"Database error: {err}")  
      return None  
def save_assessment_result(username, qb_id, score):  
   connection = create_connection()  
   if connection is None:  
      return False  
  
   try:  
      cursor = connection.cursor()  
      query = "INSERT INTO assessments (username, question_bank_id, score) VALUES (%s, %s, %s)"  
      cursor.execute(query, (username, qb_id, score))  
      connection.commit()  
      cursor.close()  
      connection.close()  
      return True  
   except mysql.connector.Error as err:  
      st.error(f"Database error: {err}")  
      return False 
def get_previous_learning_plan_end_date(username):   
  connection = create_connection()   
  if connection is None:   
    return None   
   
  cursor = connection.cursor(dictionary=True)   
  query = "SELECT end_date FROM learning_plans WHERE username = %s AND status = 'Completed' ORDER BY end_date DESC LIMIT 1"   
  cursor.execute(query, (username,))   
  result = cursor.fetchone()   
  cursor.close()   
  connection.close()   
   
  if result:   
    return datetime.strptime(result['end_date'], '%Y-%m-%d')   
  else:   
    return datetime.now()

def get_completed_date(username):   
  connection = create_connection()   
  if connection is None:   
    return None   
   
  cursor = connection.cursor(dictionary=True)   
  query = "SHOW COLUMNS FROM learning_plans LIKE 'updated_at'"   
  cursor.execute(query)   
  result = cursor.fetchone()   
  if not result:   
    query = "ALTER TABLE learning_plans ADD COLUMN updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP ON UPDATE CURRENT_TIMESTAMP"   
    cursor.execute(query)   
    connection.commit()   
   
  query = "SELECT updated_at FROM learning_plans WHERE username = %s AND status = 'Completed' ORDER BY updated_at DESC LIMIT 1"   
  cursor.execute(query, (username,))   
  result = cursor.fetchone()   
  cursor.close()   
  connection.close()   
   
  if result:   
    return datetime.strptime(result['updated_at'], '%Y-%m-%d %H:%M:%S')   
  else:   
    return datetime.now()
   
def calculate_estimated_time(num_questions, difficulty):   
  if difficulty == 'Easy':   
    estimated_time = num_questions * 10  # 10 minutes per question   
  elif difficulty == 'Medium':   
    estimated_time = num_questions * 20  # 20 minutes per question   
  elif difficulty == 'Hard':   
    estimated_time = num_questions * 30  # 30 minutes per question   
  return estimated_time
def calculate_estimated_end_date(topics, estimated_time, start_date):   
  # Calculate the estimated end date based on the topics length and other analysis   
  # For example, assume each topic takes 1 day to complete   
  num_days = len(topics)   
  estimated_end_date = (datetime.strptime(start_date, '%Y-%m-%d') + timedelta(days=num_days)).strftime('%Y-%m-%d')   
  return estimated_end_date
   
def main():
   st.title("Automated Question Builder")

   if 'user' not in st.session_state:
      st.session_state.user = None

   if st.session_state.user is None:
      tab1, tab2 = st.tabs(["Login", "Register"])

      with tab1:
         username = st.text_input("Username", key="login_username")
         password = st.text_input("Password", type="password", key="login_password")
         if st.button("Login", key="login_button"):
               user = login_user(username, password)
               if user:
                  st.session_state.user = user
                  st.success("Logged in successfully!")
                  st.rerun()
               else:
                  st.error("Invalid username or password")

      with tab2:
         new_email = st.text_input("Email", key="register_email")
         new_username = st.text_input("Username", key="register_username")
         new_password = st.text_input("Password", type="password", key="register_password")
         role = st.selectbox("Role", ["Administrator", "Trainer", "Employee"], key="register_role")
         if st.button("Register", key="register_button"):  
            if register_user(new_email, new_username, new_password, role):  
               st.success("Registration successful! Please log in.")  
            else:  
               st.error("Registration failed. Username may already exist.")  
  
   else:  
      st.sidebar.write(f"Logged in as: {st.session_state.user['username']}")  
      if st.sidebar.button("Logout", key="logout_button"):  
        st.session_state.user = None  
        st.rerun()  
  
      if st.session_state.user['role'] == 'Administrator':  
        admin_dashboard()  
      elif st.session_state.user['role'] == 'Trainer':  
        trainer_dashboard()  
      elif st.session_state.user['role'] == 'Employee':  
        employee_dashboard(st.session_state.user['username'])  
  
if __name__ == "__main__":  
   main()


